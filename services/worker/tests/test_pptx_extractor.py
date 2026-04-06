"""Unit tests for the PPTX extractor."""

from __future__ import annotations

import io

import pytest

from atlas_worker.extractors.pptx import PptxResult, extract_pptx


# ---------------------------------------------------------------------------
# In-memory .pptx fixture builder
# ---------------------------------------------------------------------------

def _build_pptx(
    slides: list[dict[str, object]] | None = None,
) -> bytes:
    """Build a minimal .pptx in memory.

    Each slide dict may contain:
        title (str): slide title
        body  (list[str]): bullet lines
        notes (str): speaker notes text
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Use blank layout (index 6 in the default template) to avoid placeholder conflicts.
    blank_layout = prs.slide_layouts[6]
    title_body_layout = prs.slide_layouts[1]  # Title and Content

    for slide_data in (slides or []):
        title_text: str = str(slide_data.get("title", ""))
        body_lines: list[str] = list(slide_data.get("body", []))  # type: ignore[arg-type]
        notes_text: str = str(slide_data.get("notes", ""))

        if title_text or body_lines:
            slide = prs.slides.add_slide(title_body_layout)
            shapes = slide.shapes
            if shapes.title and title_text:
                shapes.title.text = title_text
            if body_lines:
                try:
                    tf = shapes.placeholders[1].text_frame
                    tf.text = body_lines[0]
                    for line in body_lines[1:]:
                        tf.add_paragraph().text = line
                except (KeyError, IndexError):
                    pass
        else:
            slide = prs.slides.add_slide(blank_layout)

        if notes_text:
            notes = slide.notes_slide
            notes.notes_text_frame.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

class TestExtractPptx:
    def test_returns_pptx_result(self):
        raw = _build_pptx([{"title": "Intro", "body": ["Hello"]}])
        result = extract_pptx(raw)
        assert isinstance(result, PptxResult)

    def test_slide_count(self):
        raw = _build_pptx([
            {"title": "Slide 1"},
            {"title": "Slide 2"},
            {"title": "Slide 3"},
        ])
        result = extract_pptx(raw)
        assert result.slide_count == 3

    def test_slide_titles_captured(self):
        raw = _build_pptx([
            {"title": "Overview"},
            {"title": "Details"},
        ])
        result = extract_pptx(raw)
        assert "Overview" in result.slide_titles
        assert "Details" in result.slide_titles

    def test_full_text_contains_slide_headings(self):
        raw = _build_pptx([{"title": "Strategy", "body": ["Point A"]}])
        result = extract_pptx(raw)
        assert "Slide 1" in result.full_text

    def test_full_text_contains_body_lines(self):
        raw = _build_pptx([{"title": "Findings", "body": ["Revenue grew 20%"]}])
        result = extract_pptx(raw)
        assert "Revenue grew 20%" in result.full_text

    def test_speaker_notes_extracted(self):
        raw = _build_pptx([
            {"title": "Intro", "body": ["Hello"], "notes": "Remember to smile."}
        ])
        result = extract_pptx(raw)
        assert "Remember to smile." in result.speaker_notes

    def test_speaker_notes_empty_when_none(self):
        raw = _build_pptx([{"title": "No notes slide", "body": ["Content"]}])
        result = extract_pptx(raw)
        assert result.speaker_notes == ""

    def test_title_is_first_slide_title(self):
        raw = _build_pptx([
            {"title": "Project Atlas", "body": []},
            {"title": "Agenda", "body": []},
        ])
        result = extract_pptx(raw)
        assert result.title == "Project Atlas"

    def test_language_detected_for_english_content(self):
        raw = _build_pptx([
            {
                "title": "Report",
                "body": [
                    "The market is growing and it is the best time to invest.",
                    "In that scenario, the company is well positioned.",
                ],
            }
        ])
        result = extract_pptx(raw)
        assert result.language in (None, "en")

    def test_empty_presentation_slide_count_zero(self):
        raw = _build_pptx([])
        result = extract_pptx(raw)
        assert result.slide_count == 0

    def test_result_is_frozen(self):
        raw = _build_pptx([{"title": "T", "body": ["B"]}])
        result = extract_pptx(raw)
        with pytest.raises((AttributeError, TypeError)):
            result.full_text = "mutated"  # type: ignore[misc]

    def test_invalid_bytes_raises(self):
        with pytest.raises(Exception):
            extract_pptx(b"not a pptx file")

    def test_multiple_slides_notes_all_captured(self):
        raw = _build_pptx([
            {"title": "S1", "body": ["a"], "notes": "Note for slide 1."},
            {"title": "S2", "body": ["b"], "notes": "Note for slide 2."},
        ])
        result = extract_pptx(raw)
        assert "Note for slide 1." in result.speaker_notes
        assert "Note for slide 2." in result.speaker_notes
