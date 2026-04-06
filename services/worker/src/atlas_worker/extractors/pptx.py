"""PPTX extractor — uses python-pptx to extract slide text and speaker notes."""

from __future__ import annotations

import io
import re
from dataclasses import dataclass


@dataclass(frozen=True)
class PptxResult:
    full_text: str
    title: str | None
    slide_count: int
    slide_titles: tuple[str, ...]
    speaker_notes: str
    language: str | None


def extract_pptx(raw: bytes) -> PptxResult:
    """Extract text from every slide and speaker notes from a .pptx byte payload.

    Inputs:  raw bytes of a .pptx file.
    Outputs: PptxResult with full Markdown-formatted text, presentation title,
             slide count, per-slide titles, concatenated speaker notes, and
             detected language.
    Failures: raises ValueError if python-pptx cannot open the bytes;
              ImportError if python-pptx is not installed.
    """
    from pptx import Presentation  # deferred import
    from pptx.util import Pt  # noqa: F401 — deferred for optional dep isolation

    prs = Presentation(io.BytesIO(raw))

    slide_sections: list[str] = []
    slide_titles: list[str] = []
    notes_parts: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title = _slide_title(slide)
        if slide_title:
            slide_titles.append(slide_title)

        body_lines = _slide_body_lines(slide, skip_title=bool(slide_title))
        heading = f"## Slide {idx}" + (f": {slide_title}" if slide_title else "")
        parts = [heading]
        if body_lines:
            parts.append("\n".join(body_lines))

        notes_text = _slide_notes(slide)
        if notes_text:
            notes_parts.append(f"### Slide {idx} notes\n\n{notes_text}")

        slide_sections.append("\n\n".join(parts))

    full_text = "\n\n".join(slide_sections)
    speaker_notes = "\n\n".join(notes_parts)
    presentation_title = slide_titles[0] if slide_titles else _infer_title(full_text)
    language = _detect_language_hint(full_text)

    return PptxResult(
        full_text=full_text,
        title=presentation_title,
        slide_count=len(prs.slides),
        slide_titles=tuple(slide_titles),
        speaker_notes=speaker_notes,
        language=language,
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _slide_title(slide: object) -> str | None:
    """Return the title placeholder text from a slide, or None."""
    try:
        shapes = slide.shapes  # type: ignore[union-attr]
    except AttributeError:
        return None
    for shape in shapes:
        if shape.has_text_frame and _is_title_placeholder(shape):
            text = shape.text_frame.text.strip()
            if text:
                return text[:200]
    return None


def _is_title_placeholder(shape: object) -> bool:
    """Return True if the shape is a title or center-title placeholder."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        ph = shape.placeholder_format  # type: ignore[union-attr]
        if ph is None:
            return False
        return ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        return False


def _slide_body_lines(slide: object, *, skip_title: bool) -> list[str]:
    """Return non-empty text lines from all non-title text frames on a slide."""
    lines: list[str] = []
    try:
        shapes = slide.shapes  # type: ignore[union-attr]
    except AttributeError:
        return lines
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        if skip_title and _is_title_placeholder(shape):
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)
    return lines


def _slide_notes(slide: object) -> str:
    """Return speaker notes text from a slide, or empty string."""
    try:
        notes_slide = slide.notes_slide  # type: ignore[union-attr]
        if notes_slide is None:
            return ""
        tf = notes_slide.notes_text_frame
        if tf is None:
            return ""
        return tf.text.strip()
    except Exception:
        return ""


_LANG_HINTS: dict[str, frozenset[str]] = {
    "en": frozenset({"the", "and", "of", "to", "in", "is", "it", "that"}),
    "es": frozenset({"de", "la", "el", "en", "que", "y", "los", "las", "un"}),
    "fr": frozenset({"de", "la", "le", "les", "et", "en", "un", "une", "du"}),
    "de": frozenset({"der", "die", "das", "und", "in", "zu", "den", "ist"}),
}


def _detect_language_hint(text: str) -> str | None:
    """Return ISO-639-1 code for the most likely language, or None if uncertain."""
    words = re.findall(r"\b[a-zA-Z]{2,}\b", text.lower())
    if not words:
        return None
    sample = set(words[:500])
    scores = {lang: len(sample & hints) for lang, hints in _LANG_HINTS.items()}
    best_lang, best_score = max(scores.items(), key=lambda x: x[1])
    return best_lang if best_score >= 3 else None


def _infer_title(text: str) -> str | None:
    """Return the first heading line stripped of Markdown markers."""
    for line in text.splitlines():
        stripped = line.lstrip("#").strip()
        if stripped:
            return stripped[:200] if len(stripped) <= 200 else None
    return None
